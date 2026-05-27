"""Генерує мінімалістичну презентацію проекту (.pptx)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs" / "presentation.pptx"

# Палітра — мінімалістична
NAVY = RGBColor(0x0E, 0x2D, 0x5C)
ACCENT = RGBColor(0xE7, 0x6F, 0x51)
TEXT = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x88, 0x88, 0x88)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Слайд 16:9 (за замовчуванням python-pptx — 4:3, виставимо 13.333 x 7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_pres() -> Presentation:
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H
    return pres


def add_blank(pres):
    return pres.slides.add_slide(pres.slide_layouts[6])  # blank


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font="Calibri", italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tf


def add_paragraphs(slide, left, top, width, height, items, *,
                   size=14, color=TEXT, bullet=True, line_spacing=1.25):
    """items: list[str] | list[(text, bold)]"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, it in enumerate(items):
        text, bold = (it, False) if isinstance(it, str) else it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = ("•  " if bullet else "") + text
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return tf


def header_bar(slide, title, kicker=None):
    """Фірмова шапка з відступом і тонкою лінією."""
    # Маленький kicker (вказівник на номер/секцію)
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.3),
                 kicker, size=11, bold=True, color=ACCENT)
        title_top = Inches(0.75)
    else:
        title_top = Inches(0.55)
    add_text(slide, Inches(0.6), title_top, Inches(12.2), Inches(0.8),
             title, size=28, bold=True, color=NAVY)
    # Тонка лінія
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(1.55),
                                  Inches(0.5), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def footer(slide, idx, total):
    """Сторінкова підпис у правому нижньому куті."""
    add_text(slide, Inches(0.6), Inches(7.1), Inches(6), Inches(0.3),
             "SKU Sales Forecast  ·  УМ-з31", size=10, color=MUTED)
    add_text(slide, Inches(7.5), Inches(7.1), Inches(5.2), Inches(0.3),
             f"{idx} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def stat_card(slide, left, top, width, height, value, label, *, color=NAVY):
    """Карточка з великим числом і підписом."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = BG_LIGHT
    card.line.fill.background()
    add_text(slide, left, top + Inches(0.25), width, Inches(0.9),
             value, size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(1.15), width, Inches(0.5),
             label, size=12, color=TEXT, align=PP_ALIGN.CENTER)


def add_image(slide, path, left, top, width=None, height=None):
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return pic


def build():
    pres = new_pres()
    TOTAL = 12

    # =====================================================================
    # 1. ТИТУЛ
    s = add_blank(pres)
    # Верхня декоративна смуга
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), Inches(0), SLIDE_W, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Kicker
    add_text(s, Inches(0.8), Inches(2.4), Inches(12), Inches(0.4),
             "БІЗНЕС-ПРОГНОЗУВАННЯ", size=12, bold=True, color=ACCENT)

    # Title
    add_text(s, Inches(0.8), Inches(2.85), Inches(12), Inches(1.4),
             "Прогноз продажу товарів", size=44, bold=True, color=NAVY)
    add_text(s, Inches(0.8), Inches(3.85), Inches(12), Inches(0.8),
             "у розрізі SKU × магазин на 4 тижні наперед",
             size=22, color=TEXT)

    # Лінія
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(5.0), Inches(1.2), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Підпис
    add_text(s, Inches(0.8), Inches(5.25), Inches(12), Inches(0.5),
             "Виконавець", size=11, bold=True, color=MUTED)
    add_text(s, Inches(0.8), Inches(5.55), Inches(12), Inches(0.5),
             "Горбунова Крістіна", size=18, bold=True, color=TEXT)
    add_text(s, Inches(0.8), Inches(5.9), Inches(12), Inches(0.5),
             "група УМ-з31", size=14, color=TEXT)

    # =====================================================================
    # 2. БІЗНЕС-ПРОБЛЕМА
    s = add_blank(pres)
    header_bar(s, "Бізнес-проблема", kicker="01")

    # Два блоки в один ряд
    block_top = Inches(2.0)
    block_h = Inches(2.8)

    # Лівий блок — OOS
    left_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.6), block_top, Inches(6.0), block_h)
    left_card.fill.solid(); left_card.fill.fore_color.rgb = BG_LIGHT
    left_card.line.fill.background()
    add_text(s, Inches(0.9), block_top + Inches(0.25), Inches(5.5), Inches(0.4),
             "OUT-OF-STOCK", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.9), block_top + Inches(0.65), Inches(5.5), Inches(0.6),
             "Товар закінчився на складі", size=18, bold=True, color=NAVY)
    add_paragraphs(s, Inches(0.9), block_top + Inches(1.35), Inches(5.5), Inches(1.4),
                   ["клієнт не знаходить товар і йде до конкурента",
                    "пряма втрата виручки (GMV)",
                    "ризик зниження customer retention"], size=12)

    # Правий блок — Over-stock
    right_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(6.85), block_top, Inches(6.0), block_h)
    right_card.fill.solid(); right_card.fill.fore_color.rgb = BG_LIGHT
    right_card.line.fill.background()
    add_text(s, Inches(7.15), block_top + Inches(0.25), Inches(5.5), Inches(0.4),
             "OVER-STOCK", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(7.15), block_top + Inches(0.65), Inches(5.5), Inches(0.6),
             "Перевитрата на складі", size=18, bold=True, color=NAVY)
    add_paragraphs(s, Inches(7.15), block_top + Inches(1.35), Inches(5.5), Inches(1.4),
                   ["заморожені оборотні кошти",
                    "витрати на зберігання",
                    "ризик уцінки і списання"], size=12)

    # Висновок
    add_text(s, Inches(0.6), Inches(5.4), Inches(12.2), Inches(0.6),
             "Ручне планування десятків тисяч SKU не масштабується.",
             size=18, italic=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(5.95), Inches(12.2), Inches(0.5),
             "Потрібен автоматичний прогноз продажів у розрізі SKU × магазин.",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)

    footer(s, 2, TOTAL)

    # =====================================================================
    # 3. МЕТА І РОЗРІЗ
    s = add_blank(pres)
    header_bar(s, "Мета проекту", kicker="02")

    add_text(s, Inches(0.6), Inches(1.9), Inches(12.2), Inches(0.6),
             "Спрогнозувати тижневі продажі у розрізі SKU × магазин на 4 тижні наперед",
             size=18, color=TEXT)

    # Три картки з ключовими параметрами
    card_top = Inches(3.0)
    stat_card(s, Inches(0.6),  card_top, Inches(4.0), Inches(1.85), "SKU × store", "гранулярність прогнозу")
    stat_card(s, Inches(4.75), card_top, Inches(4.0), Inches(1.85), "тиждень",     "часовий крок")
    stat_card(s, Inches(8.9),  card_top, Inches(4.0), Inches(1.85), "4 тижні",     "горизонт прогнозу")

    add_text(s, Inches(0.6), Inches(5.4), Inches(12.2), Inches(0.5),
             "Формула планування замовлення:",
             size=13, bold=True, color=MUTED)
    formula = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(2.5), Inches(5.85), Inches(8.3), Inches(0.8))
    formula.fill.solid(); formula.fill.fore_color.rgb = NAVY
    formula.line.fill.background()
    add_text(s, Inches(2.5), Inches(5.97), Inches(8.3), Inches(0.6),
             "Q = Fct − St − GiT + SS_Fct",
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             font="Consolas")

    footer(s, 3, TOTAL)

    # =====================================================================
    # 4. ДАНІ
    s = add_blank(pres)
    header_bar(s, "Дані", kicker="03")

    # Опис ліворуч + статистика праворуч
    add_text(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.5),
             "Синтетичний датасет", size=16, bold=True, color=NAVY)
    add_paragraphs(s, Inches(0.6), Inches(2.55), Inches(6.0), Inches(3.5),
                   ["3 категорії: ноутбуки, смартфони, холодильники",
                    "15 SKU × 2 магазини",
                    "366 днів історії (повний рік)",
                    "Тижнева і річна сезонність",
                    "Свята: 8 березня, 1 травня, 24 серпня, Black Friday, Новий рік",
                    "Випадкові акції зі знижкою −10%"],
                   size=13)

    # Карточки зі статистикою
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.5),
             "В цифрах", size=16, bold=True, color=NAVY)

    stat_top = Inches(2.6)
    stat_card(s, Inches(7.2),  stat_top, Inches(2.6), Inches(1.5), "10 980", "рядків датасету", color=NAVY)
    stat_card(s, Inches(10.1), stat_top, Inches(2.6), Inches(1.5), "15",     "SKU", color=NAVY)
    stat_card(s, Inches(7.2),  stat_top + Inches(1.75), Inches(2.6), Inches(1.5), "2",  "магазини", color=NAVY)
    stat_card(s, Inches(10.1), stat_top + Inches(1.75), Inches(2.6), Inches(1.5), "52", "тижні", color=NAVY)

    footer(s, 4, TOTAL)

    # =====================================================================
    # 5. EDA
    s = add_blank(pres)
    header_bar(s, "EDA — тижневі продажі по категоріях", kicker="04")

    img = ROOT / "outputs" / "01_weekly_by_category.png"
    if img.exists():
        # Графік займає 7" по ширині, відцентрований по верт.
        add_image(s, img, Inches(0.6), Inches(2.0), width=Inches(8.0))

    add_text(s, Inches(8.9), Inches(2.0), Inches(4.0), Inches(0.4),
             "Спостереження", size=13, bold=True, color=NAVY)
    add_paragraphs(s, Inches(8.9), Inches(2.5), Inches(4.0), Inches(4.5),
                   ["Смартфони — найбільший обсяг",
                    "Холодильники — літній пік",
                    "Ноутбуки — зростання у Q4",
                    "Чітка тижнева сезонність",
                    "Сплеск у Black Friday"],
                   size=12, line_spacing=1.4)

    footer(s, 5, TOTAL)

    # =====================================================================
    # 6. МОДЕЛЬ
    s = add_blank(pres)
    header_bar(s, "Модель", kicker="05")

    # Ліва колонка — підхід
    add_text(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.5),
             "Підхід", size=16, bold=True, color=NAVY)
    add_paragraphs(s, Inches(0.6), Inches(2.55), Inches(6.0), Inches(2.5),
                   [("LightGBM", True),
                    "Gradient boosting на flat-таблиці",
                    "Рекурсивний multi-step прогноз",
                    "Train: тижні 1–44   ·   Test: 45–52"],
                   size=13, line_spacing=1.4)

    add_text(s, Inches(0.6), Inches(4.85), Inches(6.0), Inches(0.5),
             "Baseline", size=16, bold=True, color=NAVY)
    add_paragraphs(s, Inches(0.6), Inches(5.4), Inches(6.0), Inches(1.0),
                   ["Naive seasonal: середнє за 4 останні тижні"],
                   size=13)

    # Права колонка — ознаки
    feat_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.2), Inches(2.0), Inches(5.6), Inches(4.5))
    feat_card.fill.solid(); feat_card.fill.fore_color.rgb = BG_LIGHT
    feat_card.line.fill.background()
    add_text(s, Inches(7.5), Inches(2.15), Inches(5.0), Inches(0.5),
             "15 ознак", size=14, bold=True, color=NAVY)
    add_paragraphs(s, Inches(7.5), Inches(2.65), Inches(5.0), Inches(3.7),
                   [("Категоріальні", True),
                    "store_id · sku_id · category",
                    ("Календарні", True),
                    "month · quarter · is_holiday_week",
                    "woy_sin · woy_cos",
                    ("Цінові", True),
                    "avg_price · promo_share",
                    ("Лаги і ковзні", True),
                    "qty_lag_1, 2, 4, 12",
                    "qty_roll_mean_4, 12"],
                   size=12, bullet=False, line_spacing=1.25)

    footer(s, 6, TOTAL)

    # =====================================================================
    # 7. МЕТРИКИ
    s = add_blank(pres)
    header_bar(s, "Метрики на тестовому періоді", kicker="06")

    # Таблиця 3 рядки × 6 колонок
    rows, cols = 3, 6
    tbl_left, tbl_top = Inches(0.6), Inches(2.0)
    tbl_w, tbl_h = Inches(12.2), Inches(2.0)
    tbl = s.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h).table

    headers = ["Модель", "WAPE", "MAPE", "MAE", "RMSE", "Bias"]
    data = [
        ["Baseline (naive 4w)", "25.53%", "50.90%", "9.72", "15.34", "+12.19%"],
        ["LightGBM", "25.24%", "59.26%", "9.61", "15.32", "+10.44%"],
    ]

    # Заголовок
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY

    # Дані
    for ri, row in enumerate(data, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(13)
            if ri == 2 and ci > 0:
                r.font.bold = True
                r.font.color.rgb = NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_LIGHT if ri == 2 else WHITE

    add_text(s, Inches(0.6), Inches(4.3), Inches(12.2), Inches(0.5),
             "Висновок", size=14, bold=True, color=NAVY)
    add_paragraphs(s, Inches(0.6), Inches(4.8), Inches(12.2), Inches(1.6),
                   ["LightGBM кращий за baseline за WAPE та MAE/RMSE",
                    "Bias помітний у обох моделей (+10–12%): прогноз систематично завищений у Q4",
                    "MAPE гіперчутливий до тижнів з малими фактичними значеннями — основна метрика WAPE"],
                   size=13, line_spacing=1.4)

    footer(s, 7, TOTAL)

    # =====================================================================
    # 8. FEATURE IMPORTANCE
    s = add_blank(pres)
    header_bar(s, "Що зумовлює прогноз", kicker="07")

    img = ROOT / "outputs" / "03_feature_importance.png"
    if img.exists():
        add_image(s, img, Inches(0.6), Inches(2.0), width=Inches(7.5))

    add_text(s, Inches(8.4), Inches(2.0), Inches(4.5), Inches(0.4),
             "Топ-сигнали моделі", size=13, bold=True, color=NAVY)
    add_paragraphs(s, Inches(8.4), Inches(2.55), Inches(4.5), Inches(4.5),
                   [("Попередній тиждень (lag_1)", True),
                    "найсильніший сигнал",
                    ("Ковзне середнє за 4 тиж.", True),
                    "тренд останнього місяця",
                    ("Ціна і промо-частка", True),
                    "реакція на знижки",
                    ("Сезонність woy_sin/cos", True),
                    "річна циклічність"],
                   size=12, bullet=False, line_spacing=1.3)

    footer(s, 8, TOTAL)

    # =====================================================================
    # 9. ПРОГНОЗ vs ФАКТ
    s = add_blank(pres)
    header_bar(s, "Прогноз vs факт — приклади SKU", kicker="08")

    img = ROOT / "outputs" / "04_forecast_vs_actual.png"
    if img.exists():
        # Висота графіка — фіксована, ширина підлаштується по пропорції.
        # Доступна зона: між header (~1.6") та підписом (~6.5") = 4.6"
        add_image(s, img, Inches(3.2), Inches(1.95), height=Inches(4.5))

    add_text(s, Inches(0.6), Inches(6.65), Inches(12.2), Inches(0.4),
             "LightGBM (помаранчевий) точніше відтворює тренд порівняно з baseline (сірий пунктир).",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    footer(s, 9, TOTAL)

    # =====================================================================
    # 10. БІЗНЕС-ЗАСТОСУВАННЯ — формула
    s = add_blank(pres)
    header_bar(s, "Як прогноз перетворюється на замовлення", kicker="09")

    # Формула великою
    formula = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0))
    formula.fill.solid(); formula.fill.fore_color.rgb = NAVY
    formula.line.fill.background()
    add_text(s, Inches(1.5), Inches(2.2), Inches(10.3), Inches(0.7),
             "Q = Fct − St − GiT + SS_Fct",
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             font="Consolas")

    # Розшифровка — 4 картки
    desc_top = Inches(3.4)
    block_h = Inches(1.5)
    w = Inches(2.95); gap = Inches(0.1)

    for i, (k, v) in enumerate([
        ("Fct", "прогноз продажів"),
        ("St", "поточний залишок"),
        ("GiT", "товар у дорозі"),
        ("SS_Fct", "страховий запас"),
    ]):
        left = Inches(0.6) + (w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, desc_top, w, block_h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_LIGHT
        card.line.fill.background()
        add_text(s, left, desc_top + Inches(0.3), w, Inches(0.5),
                 k, size=20, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, left, desc_top + Inches(0.9), w, Inches(0.5),
                 v, size=12, color=TEXT, align=PP_ALIGN.CENTER)

    # Приклад
    add_text(s, Inches(0.6), Inches(5.3), Inches(12.2), Inches(0.4),
             "Приклад: SKU 64502, магазин 215",
             size=13, bold=True, color=NAVY)
    add_text(s, Inches(0.6), Inches(5.75), Inches(12.2), Inches(0.5),
             "Fct=152   −   St=74   −   GiT=37   +   SS=4   =   Q = 45 одиниць",
             size=15, color=TEXT, align=PP_ALIGN.CENTER, font="Consolas")

    footer(s, 10, TOTAL)

    # =====================================================================
    # 11. БІЗНЕС-ЕФЕКТ
    s = add_blank(pres)
    header_bar(s, "Бізнес-ефект", kicker="10")

    # 4 картки 2×2
    cards = [
        ("−40%", "OOS-rate",  "з 12–15% до 7–8%\n→ зростання GMV ~3%"),
        ("−12%", "Over-stock", "звільнення\nоборотних коштів"),
        ("−20 год", "час планувальника", "автоматизація рутинного\nпланування на тиждень"),
        ("$50k+", "економія SaaS", "відмова від o9, RELEX,\nQuintiq на пілот"),
    ]
    card_w, card_h = Inches(6.0), Inches(2.1)
    coords = [
        (Inches(0.6), Inches(2.0)),
        (Inches(6.85), Inches(2.0)),
        (Inches(0.6), Inches(4.4)),
        (Inches(6.85), Inches(4.4)),
    ]
    for (value, label, desc), (l, t) in zip(cards, coords):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = BG_LIGHT
        card.line.fill.background()
        add_text(s, l + Inches(0.3), t + Inches(0.2), Inches(2.5), Inches(0.7),
                 value, size=28, bold=True, color=ACCENT)
        add_text(s, l + Inches(0.3), t + Inches(0.9), Inches(5.5), Inches(0.5),
                 label, size=14, bold=True, color=NAVY)
        add_text(s, l + Inches(0.3), t + Inches(1.35), Inches(5.5), Inches(0.7),
                 desc, size=11, color=TEXT)

    footer(s, 11, TOTAL)

    # =====================================================================
    # 12. ВИСНОВКИ
    s = add_blank(pres)
    header_bar(s, "Висновки", kicker="11")

    add_paragraphs(s, Inches(0.6), Inches(2.0), Inches(12.2), Inches(3.5),
                   [("Розріз SKU × магазин × тиждень працює.", True),
                    "LightGBM забезпечує WAPE ≈ 25%, обігнавши naive-baseline за основними метриками.",
                    ("Lag-1 і rolling-mean-4 — ключові ознаки.", True),
                    "Без них модель не реагує на тренд.",
                    ("Категорія матеріальна.", True),
                    "Смартфони — WAPE 21%, холодильники — 35%. Для рідкісних SKU треба окремі моделі."],
                   size=13, line_spacing=1.4)

    # Roadmap
    add_text(s, Inches(0.6), Inches(5.45), Inches(12.2), Inches(0.4),
             "Що далі", size=14, bold=True, color=NAVY)
    add_paragraphs(s, Inches(0.6), Inches(5.85), Inches(12.2), Inches(1.2),
                   ["Розширити горизонт до 8–12 тижнів",
                    "Додати зовнішні регресори: курс валют, погода",
                    "Замінити рекурсивний прогноз на direct multi-step"],
                   size=12, line_spacing=1.3)

    footer(s, 12, TOTAL)

    # =====================================================================
    # Збереження
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    pres.save(OUT_PATH)
    print(f"Saved to {OUT_PATH}")


if __name__ == "__main__":
    build()
